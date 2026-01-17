import streamlit as st
import google.generativeai as genai
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import os
import json
import time
import random
import smtplib
from email.message import EmailMessage

# --- 1. CONFIGURATION & SECRETS ---
st.set_page_config(page_title="Study Buddy AI", page_icon="🎓", layout="wide")

# Attempt to load secrets safely
try:
    GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]
    # Check if email secrets exist, otherwise set to None (will disable email feature)
    if "EMAIL_ADDRESS" in st.secrets and "EMAIL_PASSWORD" in st.secrets:
        SENDER_EMAIL = st.secrets["EMAIL_ADDRESS"]
        SENDER_PASSWORD = st.secrets["EMAIL_PASSWORD"]
        EMAIL_ENABLED = True
    else:
        EMAIL_ENABLED = False
except FileNotFoundError:
    st.error("🚨 secrets.toml file missing! Create .streamlit/secrets.toml")
    st.stop()
except KeyError:
    st.error("🚨 Missing GOOGLE_API_KEY in secrets.toml")
    st.stop()

# --- 2. FILE HANDLING FUNCTIONS ---
USERS_FILE = "users.json"
HISTORY_FILE = "chat_history.json"

def load_users():
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, "r") as f: return json.load(f)
        except: return {}
    return {}

def save_users(users):
    with open(USERS_FILE, "w") as f: json.dump(users, f)

def load_history():
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r") as f: return json.load(f)
        except: return {}
    return {}

def save_history(session_name, messages):
    history = load_history()
    history[session_name] = messages
    with open(HISTORY_FILE, "w") as f: json.dump(history, f)

def delete_chat(session_name):
    history = load_history()
    if session_name in history:
        del history[session_name]
        with open(HISTORY_FILE, "w") as f: json.dump(history, f)

# --- 3. EMAIL & OTP FUNCTIONS ---
def send_otp_email(target_email, otp_code):
    if not EMAIL_ENABLED:
        return "Email secrets not configured."
    
    msg = EmailMessage()
    msg.set_content(f"Your Study Buddy Verification Code is: {otp_code}")
    msg['Subject'] = 'Study Buddy Password Reset'
    msg['From'] = SENDER_EMAIL
    msg['To'] = target_email

    try:
        server = smtplib.SMTP_SSL('smtp.gmail.com', 465)
        server.login(SENDER_EMAIL, SENDER_PASSWORD)
        server.send_message(msg)
        server.quit()
        return True
    except Exception as e:
        return str(e)

# --- 4. AI & FILE PROCESSING ---
def get_file_text(uploaded_file):
    text = ""
    try:
        file_ext = uploaded_file.name.split('.')[-1].lower()
        if file_ext == 'pdf':
            reader = PdfReader(uploaded_file)
            for page in reader.pages: text += page.extract_text()
        elif file_ext == 'docx':
            doc = Document(uploaded_file)
            for para in doc.paragraphs: text += para.text + "\n"
        elif file_ext == 'pptx':
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        elif file_ext == 'xlsx':
            df = pd.read_excel(uploaded_file)
            text = df.to_string()
        elif file_ext == 'txt':
            text = uploaded_file.getvalue().decode("utf-8")
    except Exception as e:
        st.error(f"Error reading file: {e}")
    return text

def get_gemini_response(api_key, prompt, context_text="", history=[]):
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-flash-latest')
    gemini_history = []
    for msg in history:
        role = "user" if msg["role"] == "user" else "model"
        gemini_history.append({"role": role, "parts": [msg["content"]]})
    
    chat = model.start_chat(history=gemini_history)
    full_prompt = f"Context: {context_text}\n\nQuestion: {prompt}" if context_text else prompt
    
    try:
        response = chat.send_message(full_prompt)
        return response.text
    except:
        return "⚠️ System busy or error. Please try again."

# --- 5. INITIALIZE SESSION STATE ---
if "users" not in st.session_state: st.session_state.users = load_users()
if "authenticated" not in st.session_state: st.session_state.authenticated = False
if "username" not in st.session_state: st.session_state.username = None
if "reset_stage" not in st.session_state: st.session_state.reset_stage = 0
if "generated_otp" not in st.session_state: st.session_state.generated_otp = None
if "reset_email" not in st.session_state: st.session_state.reset_email = None

# --- 6. AUTH HELPER FUNCTIONS ---
def login_user(username, password):
    users = load_users()
    if username in users:
        user_data = users[username]
        # Support legacy (string) and new (dict) format
        stored_password = user_data["password"] if isinstance(user_data, dict) else user_data
        
        if stored_password == password:
            st.session_state.authenticated = True
            st.session_state.username = username
            st.success("Logged in successfully!")
            st.rerun()
        else:
            st.error("Invalid password")
    else:
        st.error("User not found")

def signup_user(username, password, email):
    users = load_users()
    if username in users:
        st.error("Username already taken!")
        return

    # Check email uniqueness
    for u in users:
        if isinstance(users[u], dict) and users[u].get("email") == email:
            st.error("Email already registered!")
            return

    users[username] = {"password": password, "email": email}
    save_users(users)
    st.session_state.users = users
    st.success("Account created! Please log in.")

# --- 7. UI LOGIC ---

# === LOGIN PAGE (Black & Purple Theme) ===
if not st.session_state.authenticated:
    st.markdown("""
        <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600&display=swap');
        .stApp { background-color: #000000 !important; }
        
        .login-title {
            font-size: 3.5rem; font-weight: 700; text-align: center;
            margin-bottom: 20px; color: #d8b4fe;
            text-shadow: 0 0 25px rgba(168, 85, 247, 0.9);
        }
        
        /* The Magic Box Style */
        [data-testid="stColumn"]:nth-of-type(2) [data-testid="stVerticalBlock"] {
            background: rgba(255, 255, 255, 0.08); padding: 50px;
            border-radius: 20px; border: 1px solid rgba(255, 255, 255, 0.1);
            box-shadow: 0 0 30px rgba(255, 255, 255, 0.1);
            backdrop-filter: blur(10px); gap: 1rem;
        }

        .auth-header {
            color: rgba(255, 255, 255, 0.8); font-size: 1.1rem;
            text-align: center; margin: 0;
        }
        
        .stTextInput > div > div > input {
            background-color: rgba(0, 0, 0, 0.6) !important; color: white !important;
            border: 1px solid rgba(255, 255, 255, 0.2); border-radius: 10px; padding: 12px;
        }
        
        .stButton > button {
            width: 100%; background: linear-gradient(90deg, #7e22ce, #a855f7);
            color: white; border: none; padding: 14px;
            border-radius: 10px; font-weight: bold;
        }
        
        #MainMenu, footer, header {visibility: hidden;}
        </style>
    """, unsafe_allow_html=True)

    st.markdown('<h1 class="login-title">AI Study Buddy</h1>', unsafe_allow_html=True)
    col1, col2, col3 = st.columns([1, 1.5, 1])
    
    with col2:
        st.markdown('<p class="auth-header">Please access your account</p>', unsafe_allow_html=True)
        tab_login, tab_signup, tab_forgot = st.tabs(["Login", "Sign Up", "Forgot Password"])
        
        with tab_login:
            username = st.text_input("Username", key="login_user")
            password = st.text_input("Password", type="password", key="login_pass")
            if st.button("Log In"): login_user(username, password)
            
        with tab_signup:
            new_user = st.text_input("New Username", key="signup_user")
            new_email = st.text_input("Email Address", key="signup_email")
            new_pass = st.text_input("New Password", type="password", key="signup_pass")
            if st.button("Create Account"): 
                if new_user and new_email and new_pass:
                    signup_user(new_user, new_pass, new_email)
                else: st.warning("Please fill in all fields.")

        with tab_forgot:
            if not EMAIL_ENABLED:
                st.warning("⚠️ Email features are disabled. Configure secrets.toml.")
            else:
                # Stage 0: Enter Email
                if st.session_state.reset_stage == 0:
                    reset_email_input = st.text_input("Enter Registered Email", key="reset_email_in")
                    if st.button("Send OTP Code"):
                        users = load_users()
                        found_user = None
                        for u in users:
                            if isinstance(users[u], dict) and users[u].get("email") == reset_email_input:
                                found_user = u
                                break
                        
                        if found_user:
                            otp = str(random.randint(100000, 999999))
                            result = send_otp_email(reset_email_input, otp)
                            if result == True:
                                st.session_state.generated_otp = otp
                                st.session_state.reset_email = reset_email_input
                                st.session_state.reset_username = found_user
                                st.session_state.reset_stage = 1
                                st.rerun()
                            else:
                                st.error(f"Error sending email: {result}")
                        else:
                            st.error("Email not found.")
                
                # Stage 1: Verify OTP
                elif st.session_state.reset_stage == 1:
                    st.info(f"OTP sent to {st.session_state.reset_email}")
                    otp_input = st.text_input("Enter 6-Digit OTP", key="otp_check")
                    if st.button("Verify OTP"):
                        if otp_input == st.session_state.generated_otp:
                            st.session_state.reset_stage = 2
                            st.rerun()
                        else:
                            st.error("Invalid OTP.")
                    if st.button("Back"):
                        st.session_state.reset_stage = 0
                        st.rerun()
                
                # Stage 2: Reset Password
                elif st.session_state.reset_stage == 2:
                    st.success("Identity Verified!")
                    new_pw_reset = st.text_input("Enter New Password", type="password", key="new_pw_reset")
                    if st.button("Update Password"):
                        users = load_users()
                        users[st.session_state.reset_username]["password"] = new_pw_reset
                        save_users(users)
                        st.success("Password Updated! Please Login.")
                        time.sleep(2)
                        st.session_state.reset_stage = 0
                        st.rerun()

# === MAIN APP DASHBOARD ===
else:
    # Styles for Dashboard
    st.markdown("""
        <style>
        .stApp {
            background: linear-gradient(135deg, #2e1065 0%, #1e3a8a 100%);
            background-attachment: fixed; color: white;
        }
        .main-title {
            font-size: 3rem; font-weight: 700; text-align: center; margin-bottom: 20px;
            background: -webkit-linear-gradient(left, #a78bfa, #60a5fa);
            -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        }
        .stButton > button {
            background: rgba(255, 255, 255, 0.1) !important; color: white !important;
            border: 1px solid rgba(255, 255, 255, 0.2) !important;
            backdrop-filter: blur(10px);
        }
        .stTextInput > div > div > input {
            background-color: rgba(255, 255, 255, 0.05); color: white;
            border: 1px solid rgba(255, 255, 255, 0.2);
        }
        [data-testid="stSidebar"] {
            background-color: rgba(0, 0, 0, 0.2); backdrop-filter: blur(20px);
            border-right: 1px solid rgba(255, 255, 255, 0.1);
        }
        </style>
    """, unsafe_allow_html=True)

    if "raw_text" not in st.session_state: st.session_state.raw_text = ""
    if "messages" not in st.session_state: st.session_state.messages = []
    if "current_session_name" not in st.session_state: st.session_state.current_session_name = "New Chat"

    st.sidebar.markdown(f"### 👤 Hello, {st.session_state.username}!")
    if st.sidebar.button("🚪 Logout"): 
        st.session_state.authenticated = False
        st.session_state.username = None
        st.rerun()
        
    st.sidebar.divider()
    page = st.sidebar.radio("Go to", ["🏠 Dashboard", "📂 Study Room", "💬 AI Chat"])
    st.sidebar.divider()

    # --- PAGE 1: DASHBOARD ---
    if page == "🏠 Dashboard":
        st.markdown('<div class="main-title">📊 Learning Analytics</div>', unsafe_allow_html=True)
        history = load_history()
        total_chats = len(history)
        total_msgs = sum(len(m) for m in history.values())
        
        c1, c2, c3 = st.columns(3)
        c1.metric("Total Sessions", total_chats)
        c2.metric("Total Queries", total_msgs // 2)
        c3.metric("Status", "Active 🟢")
        
        st.subheader("Recent Activity")
        if total_chats > 0:
            data = {"Session": list(history.keys())[-5:], "Messages": [len(m) for m in list(history.values())[-5:]]}
            st.bar_chart(pd.DataFrame(data).set_index("Session"))
        else: st.info("No activity yet.")

    # --- PAGE 2: STUDY ROOM ---
    elif page == "📂 Study Room":
        st.markdown('<div class="main-title">📂 Study Room</div>', unsafe_allow_html=True)
        uploaded_file = st.file_uploader("Upload Notes (PDF, Word, PPT)", type=["pdf", "docx", "pptx", "txt", "xlsx"])
        
        if uploaded_file:
            if st.session_state.raw_text == "":
                with st.spinner("Processing File..."):
                    st.session_state.raw_text = get_file_text(uploaded_file)
                    st.success("File Processed!")
            
            st.divider()
            tab1, tab2, tab3 = st.tabs(["📝 Summarizer", "💡 Concept Simplifier", "❓ Interactive Quiz"])
            
            with tab1:
                if st.button("Generate Summary"):
                    with st.spinner("Summarizing..."):
                        res = get_gemini_response(GOOGLE_API_KEY, "Summarize this file.", st.session_state.raw_text)
                        st.markdown(res)

            with tab2:
                with st.form("simplifier_form"):
                    topic = st.text_input("Concept to explain:")
                    if st.form_submit_button("Explain") and topic:
                         with st.spinner("Simplifying..."):
                            res = get_gemini_response(GOOGLE_API_KEY, f"Explain {topic} simply.", st.session_state.raw_text)
                            st.markdown(res)

            with tab3:
                if st.button("Generate Question"):
                    prompt = """Generate ONE multiple choice question based on text. Format exactly: {"question": "...", "options": ["A", "B", "C", "D"], "answer": "B", "explanation": "..."}"""
                    res = get_gemini_response(GOOGLE_API_KEY, prompt, st.session_state.raw_text)
                    res = res.replace("```json", "").replace("```", "").strip()
                    try:
                        st.session_state.current_quiz = json.loads(res)
                        st.session_state.quiz_revealed = False
                    except: st.error("Quiz generation failed.")

                if "current_quiz" in st.session_state:
                    q = st.session_state.current_quiz
                    st.write(f"**Q:** {q['question']}")
                    user_choice = st.radio("Choose:", q['options'], key="quiz_radio")
                    if st.button("Check Answer"): st.session_state.quiz_revealed = True
                    if st.session_state.get("quiz_revealed"):
                        if user_choice == q['answer']: st.success("✅ Correct!")
                        else: st.error(f"❌ Correct: {q['answer']}")
                        st.info(q['explanation'])

    # --- PAGE 3: CHAT ---
    elif page == "💬 AI Chat":
        st.markdown('<div class="main-title">💬 AI Chat</div>', unsafe_allow_html=True)
        
        with st.sidebar:
            st.markdown("### History")
            if st.button("➕ New Chat"):
                st.session_state.messages = []
                st.session_state.current_session_name = "New Chat"
                st.rerun()
            
            saved_chats = load_history()
            for chat_name in reversed(list(saved_chats.keys())):
                c1, c2 = st.columns([0.8, 0.2])
                with c1:
                    if st.button(f"💬 {chat_name[:15]}..", key=f"load_{chat_name}"):
                        st.session_state.messages = saved_chats[chat_name]
                        st.session_state.current_session_name = chat_name
                        st.rerun()
                with c2:
                    if st.button("🗑️", key=f"del_{chat_name}"):
                        delete_chat(chat_name)
                        st.rerun()

        chat_container = st.container(height=500)
        for msg in st.session_state.messages:
            chat_container.chat_message(msg["role"]).markdown(msg["content"])

        if prompt := st.chat_input("Ask a question..."):
            st.session_state.messages.append({"role": "user", "content": prompt})
            chat_container.chat_message("user").markdown(prompt)
            
            with chat_container.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = get_gemini_response(GOOGLE_API_KEY, prompt, st.session_state.raw_text, st.session_state.messages[:-1])
                    st.markdown(response)
            
            st.session_state.messages.append({"role": "assistant", "content": response})
            if st.session_state.current_session_name == "New Chat":
                st.session_state.current_session_name = f"Chat: {prompt[:20]}..."
            save_history(st.session_state.current_session_name, st.session_state.messages)